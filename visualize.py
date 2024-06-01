from llava.model.builder import load_pretrained_model_geo
from llava.mm_utils import get_model_name_from_path
import torch
import plotly.graph_objs as go
from sklearn.cluster import SpectralClustering, DBSCAN, HDBSCAN
import umap.umap_ as umap
from utils import *
import pandas as pd
import numpy as np
from utils import get_prefecture_to_region
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
import os
from collections import defaultdict
import matplotlib.pyplot as plt
import random

class Visualizer:
    def __init__(self):
        pass
    
    def clustering(self):
        pass
    
    def visualize_pie(self):
        plt.rcParams['font.size'] = 18
        train3 = load_json('./playground/data/v8/train3.json')
        data_dict = defaultdict(int)
        prompt_dict = defaultdict(int)
        for d in train3:
            data_type = d['id'].split('_')[1]
            prompt_type = d['id'].split('_')[3]
            data_dict[data_type]+=1
            prompt_dict[data_type + '_'+prompt_type]+=1
        data_labels, data_nums, data_colors = [], [], []
        prompt_labels, prompt_nums, prompt_colors = [], [], []
        color_maps = [plt.cm.Reds, plt.cm.Greens, plt.cm.Blues, plt.cm.Oranges ]
        for i,data_type in enumerate(data_dict.keys()):
            data_nums.append(data_dict[data_type])
            data_labels.append(data_type)
            prompt_count = 1
            data_colors.append(color_maps[i](0.1*(7)))
            for j,prompt_type in enumerate(prompt_dict.keys()):
                if data_type in prompt_type:
                    prompt_labels.append(prompt_type)
                    prompt_nums.append(prompt_dict[prompt_type])
                    prompt_colors.append(color_maps[i](0.1*(prompt_count+1))) 
            
                    prompt_count+=1
        wedgeprops={"width":0.3, "edgecolor":'white'}
        print(prompt_nums, data_nums, prompt_labels, data_labels)
        plt.pie(data_nums, radius=1.3, labels=data_labels, startangle=90, counterclock=False, colors=data_colors, wedgeprops=wedgeprops)
        plt.pie(prompt_nums, labels=prompt_labels, startangle=90, counterclock=False, labeldistance=0.8, colors=prompt_colors, wedgeprops=wedgeprops)
        plt.savefig('./result/fig/data_pie.png')
        colors_sub = []
        labels_sub = []
        
    def _get_target_user(self, target_num, repeat_num=1):
        d_train = load_pkl('./preprocess/recommend/train_val.pkl')
        users = [user for user in list(d_train.keys())[:1000] if len(d_train[user]['review'])==target_num]
        target_user = random.sample(users, 1)[0]
        return target_user
    
    def visualize_generation_result(self,versions, target_num=4, repeat_num=1):
        #print('target user', target_user)
        ind2name = {'0': 'id', '1':'name', '2': 'name + desc', '3': 'name + review', '4':'name + review + image', '5': 'name + review + geo (graph)',
                    '6': 'name + review + geo (text)', '7': 'name + review + geo','12': 'user_profile'}
        ind2max = {1: 20000, 2: 20000, 3: 20000}
        prs = Presentation()
        for _ in range(repeat_num):
            target_user = self._get_target_user(target_num)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            prs.slide_width = Inches(17)
            prs.slide_height = Inches(20)
            title = slide.shapes.title
            title.text = f'ユーザー: {target_user}'
            #versions = [1,2,3,4,5,6]
            #versions = [1,2,3]
            rows, cols = len(versions)+1, 3
            top, left, width, height = Inches(2),Inches(0.5),Inches(20),Inches(0.8)
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            # Set column widths
            table.columns[0].width = Inches(1.5)
            #table.columns[1].width = Inches(7)
            table.columns[1].width = Inches(2)
            #table.columns[2].width = Inches(2)
            table.columns[2].width = Inches(10)
            #table.columns[4].width = Inches(5)

            # Set headers
            table.cell(0, 0).text = "Strategy"
            # table.cell(0, 1).text = "Prompt"
            # table.cell(0, 2).text = "Rec_Pred"
            # table.cell(0, 3).text = "Rec_GT"
            # table.cell(0, 4).text = "Review_Pred"
            # table.cell(0, 5).text = "Review_GT"
            table.cell(0, 1).text = "Rec_Pred"
            #table.cell(0, 2).text = "Rec_GT"
            table.cell(0, 2).text = "Review_Pred"
            #table.cell(0, 4).text = "Review_GT"
            
            datas = [load_json(f'./playground/data/v8/test{v}.json') for v in versions]
            id2inds = [{data[i]['id']:i for i in range(len(data))} for data in datas]
            rec_results = []
            for i in versions:
                if os.path.exists(f'./result/rec_direct/llava-v1.5-13b-jalan-review-lora-v8.{i}_0_20000.csv'):
                    rec_results.append(pd.read_csv(f'./result/rec_direct/llava-v1.5-13b-jalan-review-lora-v8.{i}_0_20000.csv'))
                else:rec_results.append(None)
                
            review_results = []
            for i in versions:
                max_user = ind2max.get(i, 1000)
                #print(f'./result/rec_review/llava-v1.5-13b-jalan-review-lora-v8.{i}_0_{max_user}.csv', os.path.exists(f'./result/rec_review/llava-v1.5-13b-jalan-review-lora-v8.{i}_0_{max_user}.csv'))
                if os.path.exists(f'./result/rec_review/llava-v1.5-13b-jalan-review-lora-v8.{i}_0_{max_user}.csv'):
                    review_results.append(pd.read_csv(f'./result/rec_review/llava-v1.5-13b-jalan-review-lora-v8.{i}_0_{max_user}.csv'))
                else:review_results.append(None)
            #rec_results = [pd.read_csv(f'./result/rec_direct/llava-v1.5-13b-jal an-review-lora-v7.{i}_0_20000.csv') for i in versions if os.path.exists(f'./result/rec_direct/llava-v1.5-13b-jalan-review-lora-v7.{i}_0_20000.csv')else None]
            #review_results = [pd.read_csv(f'./result/rec_reviews/user_history_llava-v1.5-13b-jalan-review-lora-v7.{i}_1000.csv') for i in versions if os.path.exists(f'./result/rec_reviews/user_history_llava-v1.5-13b-jalan-review-lora-v7.{i}_1000.csv') else None]

            for i,v in enumerate(versions):
                rec_tmp = rec_results[i]
                review_tmp = review_results[i]
                if rec_tmp is not None:
                    if 'user' in rec_tmp.columns:user_col = 'user'
                    elif 'users' in rec_tmp.columns:user_col='users'
                    #print(rec_tmp.query(f'{user_col}=="{target_user}"'))
                    try:
                        rec_tmp = rec_tmp.query(f'{user_col}=="{target_user}"').reset_index().loc[0, f'output{v}']
                        print(rec_tmp)
                    except KeyError:
                        rec_tmp = ''
                if review_tmp is not None:
                    if 'user' in review_tmp.columns:user_col = 'user'
                    elif 'users' in review_tmp.columns:user_col='users'
                    try:
                        review_tmp = review_tmp.query(f'{user_col}=="{target_user}"').reset_index().loc[0, f'output{v}']
                        print(review_tmp)
                    except KeyError:
                        review_tmp = ''
                data_review = datas[i][id2inds[i][f'{target_user}_review_{v}_1_False']]
                data_rec = datas[i][id2inds[i][f'{target_user}_direct_{v}_1']]
                #print('data', data)
                table.cell(i+1, 0).text = ind2name[str(v)]
                table.cell(i+1, 1).text = rec_tmp    
                #table.cell(i+1, 2).text = data_rec['conversations'][1]['value']
                table.cell(i+1, 2).text = review_tmp
                #table.cell(i+1, 4).text = data_review['conversations'][1]['value']
                
            top, left, width, height = Inches(2), Inches(0.5), Inches(20), Inches(0.8)
            new_top = top + height + Inches(10)  # 新しいテーブルの位置を最初のテーブルの下に配置
            new_table = slide.shapes.add_table(4, 2, left, new_top, width, Inches(2)).table 
            new_table.columns[0].width = Inches(2)
            #table.columns[1].width = Inches(7)
            new_table.columns[1].width = Inches(15)
            new_table.cell(0, 0).text = "type"
            new_table.cell(0, 1).text = ""
            new_top = top + height + Inches(1)  # 新しいテーブルの位置を最初のテーブルの下に配置
            new_table.cell(1, 0).text = "Spot (GT)"
            new_table.cell(2, 0).text = "Review (GT)"
            new_table.cell(3, 0).text = "User Reviews"
            
            i = 2 
            v = versions[i]
            new_table.cell(1, 1).text = datas[i][id2inds[i][f'{target_user}_direct_{v}_1']]['conversations'][1]['value']
            new_table.cell(2, 1).text = data_review['conversations'][1]['value']
            new_table.cell(3, 1).text = data_review['conversations'][0]['value']
                #table.cell(i+1, 1).text = data_review['conversations'][0]['value']
                # table.cell(i+1, 2).text = rec_tmp    
                # table.cell(i+1, 3).text = data_rec['conversations'][1]['value']
                # table.cell(i+1, 4).text = review_tmp
                # table.cell(i+1, 5).text = data_review['conversations'][1]['value']
                
        prs.save(f'./result/slide/case_study_{target_num}_{repeat_num}.pptx')
            
            
    
    def visualize_geo_embedding(self, forward=False):
        if forward:
            model_path = './checkpoints/llava-v1.5-13b-jalan-review-lora-v8.10'
            model_name = get_model_name_from_path(model_path)
            model_base = 'lmsys/vicuna-13b-v1.5'
            tokenizer, model, image_processor, context_len = load_pretrained_model_geo(model_path, model_base, model_name, False, False, 'cuda')
            geo_tower = model.get_model().get_geo_tower()
            embs = geo_tower.forward(torch.arange(geo_tower.num_nodes))
            print('embs', embs.shape)
            torch.save(embs, './result/emb/geo_emb8.10.pt')
        embs = torch.load('./result/emb/geo_emb8.10.pt').float()
        df_review = pd.read_csv('./preprocess/recommend/filtered_review_df.csv')
        df_exp = pd.read_pickle('/home/yamanishi/project/airport/src/data/experiment.pkl')
        ind2spot = {i:spot for i,spot in enumerate(df_review['spot'].unique())}
        spot2pref = dict(zip(df_exp['spot_name'], df_exp['prefecture']))
        spot_names = [ind2spot[i] for i in range(len(ind2spot))]
        prefs = [spot2pref.get(spot, '不明') for spot in spot_names]
        spot_names = [spot_names[i] + ' '+ prefs[i] for i in range(len(spot_names))]
        _, labels_pref = np.unique(prefs, return_inverse=True)
        plot_2d(embs, labels_pref, spot_names, save_path='./result/emb/geo_emb_pref810.html')
        pref2region = get_prefecture_to_region()
        regions = [pref2region.get(pref, '不明') for pref in prefs]
        _, labels = np.unique(regions, return_inverse=True)
        print('labels', labels)
        spot_names = [spot_names[i] + ' '+ regions[i] for i in range(len(spot_names))]
        plot_2d(embs, labels, spot_names, save_path='./result/emb/geo_emb_region810.html')

        
if __name__=='__main__':
    vis = Visualizer()
    #vis.visualize_pie()
    vis.visualize_generation_result(versions=[1,2,3, 12], target_num=4, repeat_num=10)
    exit()
    vis.visualize_geo_embedding(forward=True)
    exit()
    vis.visualize_generation_result(versions=[1,2,3, 5], target_num=4, repeat_num=10)
    exit()
    vis.visualize_geo_embedding()