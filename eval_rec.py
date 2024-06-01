import pandas as pd
from tqdm import tqdm
from tqdm.contrib import tenumerate
import json
from eval import Evaluator, AverageMeter
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
import re
import os
import argparse


class RecEvaluator(Evaluator):
    def __init__(self):
        super().__init__()
        self.df = pd.read_csv('./preprocess/recommend/filtered_review_df.csv')
        self.users = self.df['name'].unique()
        self.ind2name = {'0': 'id', '1':'name', '2': 'name + desc', '3': 'name + review', '4':'name + review + image', '5': 'name + review + geo (positional)',
            '6': 'name + review + geo (text)', '7': 'name + geo (id)',
            '8': 'name + geo (text)', '9': 'name + review + aspect loss',
            '10': 'name + review + contra loss', '11': 'PEPLER [TOIS23]',
            '12': 'user profile', '13': 'user profile + item profile','18': 'vicuna + review + item profile',
            '14': 'Name + review + item profile', '15': 'Name + review (short) + item profile', 
            '16': 'Name + user profile', '17':'Name + user profile + item profile','18':'Vicuna+ review + item profile', '19': 'GPT + review + item profile',
            '20': 'Name + user profile + item profile + matching', '21': 'Name + review + item profile + matching', '22': 'Name + Review + Matching'}
        
    @staticmethod
    def get_version(path):
        pattern = r'lora-v\d+\.(\d+)'

        # パターンにマッチする部分を検索
        match = re.search(pattern, path)

        if match:
            # マッチした部分の最初のグループを取得
            result = match.group(1)
            #print("取得した数字:", result)
        #else:
            #print("マッチするものがありませんでした。")
        return result
    
    def evaluate_review(self,df_paths):
        metrics = []
        for df_path in df_paths:
            metric1 = self._calc_review_quality_metric(df_path, 'output1', 'gt')
            metric2 = self._calc_review_quality_metric(df_path, 'output2', 'gt')
            metric3 = self._calc_review_quality_metric(df_path, 'output3', 'gt')
            metrics.append((metric1, metric2, metric3))
        print(df_paths, metrics)
        
    def _calculate_recommend_metric(self, df_path, gt_row, pred_row, topk=5, max_user=-1):
        df = pd.read_csv(df_path)
        if max_user!=-1:df = df[:max_user]
        #print(df.head())
        metrics = {}
        hr = AverageMeter()
        recommend_num = AverageMeter()
        count=0
        for i,(gt, pred) in enumerate(zip(df[gt_row], df[pred_row])):
            if pd.isna(pred):continue
            gt = gt.replace('<s> ', '')
            gt = gt.replace('</s>', '')
            pred = pred.split(',')
            pred = list(dict.fromkeys(pred))
            pred = pred[:min(len(pred), topk)]
            recommend_num.update(len(pred))
            pred = [p.replace('<s> ', '').replace('</s>', '') for p in pred]
            hr.update(int(gt in pred))
            count += int(gt in pred)
        metrics['hr'] = hr.get_average()
        print('recommend num', recommend_num.get_average())
        #print('count', count)
        return metrics
            
    def evaluate_rec(self, df_paths, k=[1,3,5], outputs=[1,2,3]):
        metrics = []
        for df_path in df_paths:
            for k in [1,3,5]:
                for output in outputs:
                    metric = self._calculate_recommend_metric(df_path, 'gt', f'output{output}', topk=k)
                    metrics.append((df_path,f'top{k}', f'prompt{output}', metric))
                
        print(metrics)
        
    def make_slide_rec(self, df_paths, k=[1,3,5], ):
        ind2name = self.ind2name
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        metrics = []
        for df_path in df_paths:
            ind = RecEvaluator.get_version(df_path)
            result = {}
            for k in [1,3,5]:
                if not os.path.exists(df_path):continue
                metric = self._calculate_recommend_metric(df_path, 'gt', f'output{ind}', topk=k)
                print('metric', metric)
                result[f'HR@{k}'] = metric['hr']
            metrics.append((df_path, result))
                
        rows, cols = len(metrics) + 1, 4
        top, left, width, height = Inches(2),Inches(0.5),Inches(9),Inches(0.8)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths
        table.columns[0].width = Inches(2.5)
        table.columns[1].width = Inches(2.2)
        table.columns[2].width = Inches(2.2)
        table.columns[3].width = Inches(2.2)

        # Set headers
        table.cell(0, 0).text = "Model"
        table.cell(0, 1).text = "HR@1"
        table.cell(0, 2).text = "HR@3"
        table.cell(0, 3).text = "HR@5"

        # Populate table rows
        for i, (model_path, result) in enumerate(metrics, 1):
            print(model_path, result)
            table.cell(i, 0).text = ind2name[RecEvaluator.get_version(model_path)]
            table.cell(i, 1).text = f"{result.get('HR@1', 0):.4f}"
            table.cell(i, 2).text = f"{result.get('HR@3', 0):.4f}"
            table.cell(i, 3).text = f"{result.get('HR@5', 0):.4f}"

            # Align text to center
            for j in range(cols):
                table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        prs.save('./result/slide/recommendation_metrics_overview.pptx')
        
    def make_slide_review(self, df_paths, max_user=-1,):
        ind2name  = self.ind2name
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        metrics = []
        for df_path in df_paths:
            ind = RecEvaluator.get_version(df_path)
            result = {}
            if not os.path.exists(df_path):continue
            if '14' in df_path:ind=14
            metric = self._calc_review_quality_metric(df_path, 'gt', f'output{ind}', max_user=max_user)
            for metric_name in ['bleu', 'rouge_1', 'rouge_2', 'rouge_l']:
                result[metric_name] = metric[metric_name]
            metrics.append((df_path, result))
                
        rows, cols = len(metrics) + 1, 5
        top, left, width, height = Inches(2),Inches(0.5),Inches(9),Inches(0.8)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths
        table.columns[0].width = Inches(2.5)
        table.columns[1].width = Inches(2.2)
        table.columns[2].width = Inches(2.2)
        table.columns[3].width = Inches(2.2)
        table.columns[4].width = Inches(2.2)

        # Set headers
        table.cell(0, 0).text = "Model"
        table.cell(0, 1).text = "BLEU"
        table.cell(0, 2).text = "ROUGE1"
        table.cell(0, 3).text = "ROUGE2"
        table.cell(0, 4).text = "ROUGEL"

        # Populate table rows
        for i, (model_path, result) in enumerate(metrics, 1):
            print(model_path, result)
            table.cell(i, 0).text = ind2name[RecEvaluator.get_version(model_path)]
            table.cell(i, 1).text = f"{result.get('bleu', 0):.4f}"
            table.cell(i, 2).text = f"{result.get('rouge_1', 0):.4f}"
            table.cell(i, 3).text = f"{result.get('rouge_2', 0):.4f}"
            table.cell(i, 4).text = f"{result.get('rouge_l', 0):.4f}"

            # Align text to center
            for j in range(cols):
                table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        prs.save('./result/slide/recommendation_metrics_review.pptx')
                    
if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--eval_type', default='rec')
    parser.add_argument('--ind', default=1)
    evaluator = RecEvaluator()
    # evaluator.make_slide_rec([f'./result/rec_direct/llava-v1.5-13b-jalan-review-lora-v8.{i}_0_20000.csv'
    #                           for i in [1,2,3,14, 16, 17]])
    # exit()
    #exit()
    evaluator.make_slide_review(['./result/rec_review/llava-v1.5-13b-jalan-review-lora-v8.20_0_1000.csv',], max_user=1000)
    #                            './result/rec_review/llava-v1.5-13b-jalan-review-lora-v8.18_0_1000.csv'], max_user=1000)
    exit()
                                # './result/rec_review/llava-v1.5-13b-jalan-review-lora-v8.3_14_0_1000.csv',
                                # './result/rec_review/llava-v1.5-13b-jalan-review-lora-v8.1_0_20000.csv',
                                #  './result/rec_review/llava-v1.5-13b-jalan-review-lora-v8.2_0_20000.csv',
                                #  './result/rec_review/llava-v1.5-13b-jalan-review-lora-v8.3_0_20000.csv',
                                #  './result/rec_review/llava-v1.5-13b-jalan-review-lora-v8.4_0_1000.csv',
                                #  './result/rec_review/llava-v1.5-13b-jalan-review-lora-v8.5_0_1000.csv',
                                #  './result/rec_review/llava-v1.5-13b-jalan-review-lora-v8.9_0_1000.csv',
                                #  './result/rec_review/llava-v1.5-13b-jalan-review-lora-v8.12_0_1000.csv'],
                                # max_user=1000)
    exit()
    evaluator.evaluate_review(['./result/rec_reviews/user_history_llava-v1.5-13b-jalan-review-lora-v7.5_1000.csv',], outputs=[5])
    evaluator.make_slide_review([f'./result/rec_reviews/user_history_llava-v1.5-13b-jalan-review-lora-v7.{ind}_1000.csv' for ind in range(1, 6)])
    #evaluator.make_slide_rec([f'./result/rec_direct/llava-v1.5-13b-jalan-review-lora-v7.{ind}_0_20000.csv' for ind in range(1, 6)])
    exit()
    ind = 2
    evaluator.evaluate_rec([f'./result/rec_direct/llava-v1.5-13b-jalan-review-lora-v7.{ind}_0_20000.csv'], outputs=[ind])
    ind = 5
    evaluator.evaluate_rec([f'./result/rec_direct/llava-v1.5-13b-jalan-review-lora-v7.{ind}_0_20000.csv'], outputs=[ind])
    exit()
    evaluator.evaluate_rec(['./result/rec_direct/llava-v1.5-13b-jalan-review-lora-v7.2_0_1000.csv',], outputs=[1,2,3])
    evaluator.evaluate_rec(['./result/rec_direct/llava-v1.5-13b-jalan-review-lora-v7.3_1000.csv',], outputs=[3])
    evaluator.evaluate_rec(['./result/rec_direct/llava-v1.5-13b-jalan-review-lora-v7.5_0_1000.csv',], outputs=[5])
    exit()
    evaluator.evaluate_review(['./result/rec_reviews/user_history_llava-v1.5-13b-jalan-review-lora-v7.5_1000.csv',], outputs=[5])
    exit()
    evaluator.evaluate_review(['./result/rec_reviews/user_history_llava-v1.5-13b-jalan-review-lora-v7.2_1000.csv',])
    exit()
    evaluator.evaluate_rec(['./result/rec_direct/llava-v1.5-13b-jalan-review-lora-v7.1_1000.csv',
                               #'./result/rec_direct/llava-v1.5-13b-jalan-review-lora-v7.2_100.csv',
                               './result/rec_direct/llava-v1.5-13b-jalan-review-lora-v7.3_1000.csv']) 
    exit()
    evaluator.evaluate_review(['./result/rec_reviews/user_history_llava-v1.5-13b-jalan-review-lora-v7.1_1000.csv',
                               './result/rec_reviews/user_history_llava-v1.5-13b-jalan-review-lora-v7.3_1000.csv'])        
    #evaluator.evaluate_review('./result/rec_reviews/user_history_llava-v1.5-13b-jalan-review-lora-v7.3_1000.csv')     